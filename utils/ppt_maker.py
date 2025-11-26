from pptx import Presentation
from pptx.util import Inches
import os


def generate_ppt(slide_paths, output_path):
    """Generate a PPTX from slide_paths and save to output_path.

    output_path should be something like jobs/<job_id>/output.pptx.
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # empty slide layout

    for slide_img in slide_paths:
        slide = prs.slides.add_slide(blank_layout)

        # Insert full-screen image
        slide.shapes.add_picture(
            slide_img,
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # Ensure parent directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs.save(output_path)

    return output_path
